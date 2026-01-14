import os
import streamlit as st
from openai import OpenAI
from pathlib import Path
from datetime import datetime
from typing import List, Dict

# 페이지 설정
st.set_page_config(
    page_title="doorong 전용 유아발달상황 평가 작성 도우미",
    page_icon="🎓",
    layout="wide"
)

BASE_DIR = Path(__file__).resolve().parent
UPLOAD_DIR = BASE_DIR / "uploads"
SAVE_DIR = BASE_DIR / "saved_evaluations"

AUTH_USERNAME = "doorong"
AUTH_PASSWORD = "1234"

# 세션 스테이트 초기화
if "authenticated" not in st.session_state:
    st.session_state.authenticated = False
if "messages" not in st.session_state:
    st.session_state.messages = []
if "rag_documents" not in st.session_state:
    st.session_state.rag_documents = []
if "documents_loaded" not in st.session_state:
    st.session_state.documents_loaded = False
if "last_rag_context" not in st.session_state:
    st.session_state.last_rag_context = ""


def require_login() -> None:
    if st.session_state.authenticated:
        return

    st.title("🔒 로그인")
    st.info("지정된 아이디/비밀번호로 로그인하세요.")

    with st.form("login_form"):
        username = st.text_input("ID")
        password = st.text_input("Password", type="password")
        submitted = st.form_submit_button("로그인")

    if submitted:
        if username == AUTH_USERNAME and password == AUTH_PASSWORD:
            st.session_state.authenticated = True
            st.success("로그인 완료")
            st.rerun()
        else:
            st.error("아이디 또는 비밀번호가 올바르지 않습니다.")

    st.stop()


def get_api_key() -> str:
    try:
        return st.secrets.get("OPENAI_API_KEY", os.environ.get("OPENAI_API_KEY", ""))
    except Exception:
        return os.environ.get("OPENAI_API_KEY", "")


def extract_text_from_uploaded_file(uploaded_file) -> str:
    """업로드된 파일에서 텍스트 추출"""
    try:
        content = uploaded_file.read()

        if uploaded_file.type == "text/plain":
            return content.decode("utf-8")
        if uploaded_file.type == "application/pdf":
            try:
                import PyPDF2
                from io import BytesIO

                pdf_reader = PyPDF2.PdfReader(BytesIO(content))
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text()
                return text
            except Exception:
                return f"[PDF 파일: {uploaded_file.name} - 내용 추출 실패]"
        if uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            try:
                from pptx import Presentation
                from io import BytesIO

                prs = Presentation(BytesIO(content))
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                return text
            except Exception:
                return f"[PPTX 파일: {uploaded_file.name} - 내용 추출 실패]"

        return f"[{uploaded_file.name}: 지원하지 않는 형식]"

    except Exception as exc:
        return f"[파일 읽기 오류: {str(exc)}]"


def create_embeddings_simple(text: str, client: OpenAI) -> List[float]:
    """OpenAI Embeddings API를 사용하여 임베딩 생성"""
    try:
        response = client.embeddings.create(
            model="text-embedding-3-small",
            input=text[:8000]
        )
        return response.data[0].embedding
    except Exception as exc:
        st.error(f"임베딩 생성 오류: {str(exc)}")
        return []


def cosine_similarity(vec1: List[float], vec2: List[float]) -> float:
    """코사인 유사도 계산"""
    if not vec1 or not vec2:
        return 0.0
    dot_product = sum(a * b for a, b in zip(vec1, vec2))
    magnitude1 = sum(a * a for a in vec1) ** 0.5
    magnitude2 = sum(b * b for b in vec2) ** 0.5
    if magnitude1 == 0 or magnitude2 == 0:
        return 0.0
    return dot_product / (magnitude1 * magnitude2)


def search_similar_documents(query: str, documents: List[Dict], client: OpenAI, top_k: int = 3) -> List[Dict]:
    """쿼리와 유사한 문서 검색"""
    if not documents:
        return []

    query_embedding = create_embeddings_simple(query, client)
    if not query_embedding:
        return []

    similarities = []
    for doc in documents:
        similarity = cosine_similarity(query_embedding, doc.get("embedding", []))
        similarities.append({
            "document": doc,
            "similarity": similarity
        })

    similarities.sort(key=lambda x: x["similarity"], reverse=True)
    return [item["document"] for item in similarities[:top_k]]


require_login()

# 사이드바 - 설정
with st.sidebar:
    st.title("⚙️ 설정")

    mode = st.radio(
        "모드 선택",
        ["기본 (누리과정만)", "고급 (문서 업로드 포함)"]
    )

    if mode == "고급 (문서 업로드 포함)":
        st.subheader("📚 참고 문서 업로드")
        st.info("평가 작성에 참고할 자료를 업로드하세요")

        uploaded_files = st.file_uploader(
            "파일 선택 (TXT, PDF, PPTX)",
            type=["txt", "pdf", "pptx"],
            accept_multiple_files=True,
            label_visibility="collapsed"
        )

        if uploaded_files and st.button("📂 문서 분석 및 로드"):
            api_key = get_api_key()
            if not api_key:
                st.error("먼저 환경 변수 또는 Streamlit Secrets에 API 키를 설정해주세요!")
            else:
                with st.spinner("문서를 분석하고 있습니다... (시간이 걸릴 수 있습니다)"):
                    try:
                        client = OpenAI(api_key=api_key)
                        st.session_state.rag_documents = []

                        UPLOAD_DIR.mkdir(exist_ok=True)

                        for uploaded_file in uploaded_files:
                            file_path = UPLOAD_DIR / uploaded_file.name
                            with open(file_path, "wb") as f:
                                f.write(uploaded_file.getbuffer())

                            text = extract_text_from_uploaded_file(uploaded_file)
                            chunks = [text[i:i + 2000] for i in range(0, len(text), 2000)]

                            for i, chunk in enumerate(chunks):
                                if chunk.strip():
                                    embedding = create_embeddings_simple(chunk, client)

                                    st.session_state.rag_documents.append({
                                        "filename": uploaded_file.name,
                                        "chunk_id": i,
                                        "content": chunk,
                                        "embedding": embedding
                                    })

                        st.session_state.documents_loaded = True
                        st.success(
                            f"✅ {len(uploaded_files)}개 파일에서 {len(st.session_state.rag_documents)}개 문서 조각을 분석했습니다!"
                        )

                    except Exception as exc:
                        st.error(f"문서 처리 중 오류 발생: {str(exc)}")

        if st.session_state.documents_loaded:
            st.success(f"📖 현재 로드된 문서: {len(st.session_state.rag_documents)}개 조각")

    st.subheader("📋 사용 가이드")
    st.markdown(
        """
    **1단계**: OpenAI API 키를 입력하세요

    **2단계**: (고급 모드) 참고 문서를 업로드하세요
    - 기존 평가 사례 (PPTX)
    - 발달 단계 자료 (PDF, TXT)
    - 평가 작성 가이드 등

    **3단계**: 아이 정보를 입력하세요
    - 이름, 나이, 관찰 내용

    **4단계**: AI가 생성한 평가를 확인하고 수정하세요
    """
    )

    if st.button("🗑️ 대화 초기화"):
        st.session_state.messages = []
        st.rerun()

    if st.button("🔓 로그아웃"):
        st.session_state.authenticated = False
        st.rerun()

# 메인 영역
st.title("🎓 doorong 전용 유아발달상황 종합평가 작성 도우미")
st.markdown("---")

# API 키 확인
api_key = get_api_key()
if not api_key:
    st.warning("⚠️ 환경 변수 또는 Streamlit Secrets에 OpenAI API 키를 설정해주세요.")
    st.info(
        """
    ### OpenAI API 키 설정 방법
    - 로컬 실행: `OPENAI_API_KEY` 환경 변수로 설정
    - Streamlit Cloud: Secrets에 `OPENAI_API_KEY` 추가
    """
    )
    st.stop()

# OpenAI 클라이언트 초기화
client = OpenAI(api_key=api_key)

# 시스템 프롬프트 정의
SYSTEM_PROMPT = """당신은 경력 10년차 유치원 교사이자 유아교육 전문가입니다.
유아발달상황 종합평가를 작성할 때 다음 원칙을 반드시 지켜주세요:

## 작성 원칙
1. **따뜻하고 상냥한 교사 말투** 사용
2. **아이의 행동 관찰 → 발달적 의미 해석 → 교사의 전문성** 순서로 작성
3. **5개 발달 영역 모두 포함**:
   - 신체운동건강
   - 의사소통
   - 사회관계
   - 예술경험
   - 자연탐구

4. **구조**:
   - 도입: 아이의 전반적인 모습과 성장
   - 본문: 각 영역별 구체적 관찰과 발달적 의미
   - 마무리: "~의 2학기 생활도 응원합니다!" + 가정 연계 지원 방안

5. **문체**:
   - 자연스러운 문맥 연결
   - 구체적 행동 사례 제시
   - 긍정적이고 격려하는 톤
   - 전문 용어를 쉽게 풀어서 설명

## 참고 자료 활용
제공된 참고 문서를 바탕으로 각 아이의 개별성을 존중하여 작성하세요.

## 출력 형식
평가문을 작성한 후, 마지막에 다음 정보를 추가하세요:

---
**📚 참고한 자료**
- [참고한 문서명 또는 발달 영역]

**🔍 발달적 근거**
- [해당 행동이 나타내는 발달 단계와 의미]
"""

# 기본 발달 지식 베이스
DEFAULT_KNOWLEDGE = """
# 누리과정 5개 영역별 발달 특성

## 1. 신체운동·건강
### 만 3세
- 대근육: 달리기, 점프하기, 공 던지고 받기 등 기본 운동 능력 발달
- 소근육: 가위질, 블록 쌓기, 그리기 도구 사용 시작
- 건강: 스스로 먹기, 손 씻기 등 기본 생활습관 형성

### 만 4세
- 대근육: 한발로 서기, 리듬에 맞춰 움직이기, 공놀이 조절력 향상
- 소근육: 선 따라 자르기, 단추 끼우기, 작은 블록 조립
- 건강: 식사 예절 알기, 규칙적인 생활습관 형성

### 만 5세
- 대근육: 줄넘기, 공 차고 받기, 평균대 건너기 등 복합 운동 가능
- 소근육: 정교한 가위질, 젓가락 사용, 글씨 쓰기 시도
- 건강: 건강한 음식 선택, 안전 규칙 이해와 실천

## 2. 의사소통
### 만 3세
- 듣기: 간단한 이야기 집중해서 듣기
- 말하기: 3-4단어 문장으로 경험 말하기
- 읽기: 그림책 관심, 반복되는 말 따라하기
- 쓰기: 끼적이기, 선 긋기

### 만 4세
- 듣기: 이야기 순서 이해하기
- 말하기: 5-6단어 복문으로 생각과 느낌 표현
- 읽기: 글자에 관심, 자신의 이름 읽기
- 쓰기: 의미 있는 끼적이기, 자기 이름 쓰기 시도

### 만 5세
- 듣기: 이야기 내용 이해하고 질문에 답하기
- 말하기: 경험을 순서대로 구체적으로 말하기
- 읽기: 친숙한 글자 읽기, 책 내용 이해
- 쓰기: 자기 이름과 친숙한 글자 쓰기

## 3. 사회관계
### 만 3세
- 자기존중: 자신의 특성과 선호 알아가기
- 정서인식: 기본 감정(기쁨, 슬픔, 화남) 표현하기
- 또래관계: 병행놀이, 짧은 시간 함께 놀이하기
- 협력: 교사의 도움으로 순서 지키기

### 만 4세
- 자기존중: 자신의 장점과 소중함 알기
- 정서조절: 감정을 말로 표현하기 시작
- 또래관계: 협동놀이 시작, 친한 친구 생기기
- 협력: 간단한 규칙 지키기, 역할 나누기

### 만 5세
- 자기존중: 자신감 있게 행동하기
- 정서조절: 감정 조절 시도, 타인 감정 이해
- 또래관계: 안정적 또래관계, 갈등 해결 시도
- 협력: 규칙의 필요성 이해, 공동 목표를 위한 협력

## 4. 예술경험
### 만 3세
- 탐색: 다양한 미술 재료 탐색
- 표현: 자유로운 끼적이기, 만들기
- 감상: 자신의 작품에 관심 갖기
- 음률: 노래 부르기 즐기기, 리듬 따라하기

### 만 4세
- 탐색: 색, 모양, 질감의 차이 인식
- 표현: 의도를 가지고 그리기, 만들기
- 감상: 친구 작품에 관심 갖고 이야기하기
- 음률: 노래에 맞춰 움직이기, 간단한 악기 연주

### 만 5세
- 탐색: 예술 요소 변별하기
- 표현: 다양한 방법으로 창의적 표현
- 감상: 작품의 특징 발견하고 설명하기
- 음률: 리듬과 노래를 창의적으로 표현

## 5. 자연탐구
### 만 3세
- 탐구과정: 주변 사물과 자연 현상에 관심
- 생명존중: 주변 동식물에 관심 갖기
- 수학적 탐구: 사물의 색, 모양 인식
- 과학적 탐구: 물체의 특성 감각으로 탐색

### 만 4세
- 탐구과정: 궁금한 것을 질문하고 알아보기
- 생명존중: 동식물 특성 관찰하고 소중히 여기기
- 수학적 탐구: 5까지 수 세기, 패턴 만들기
- 과학적 탐구: 물체와 물질의 기본 특성 알기

### 만 5세
- 탐구과정: 궁금한 것을 탐구과정을 통해 알아가기
- 생명존중: 생명체 성장과정 이해, 생명 존중
- 수학적 탐구: 10까지 수 세기, 간단한 덧셈과 뺄셈
- 과학적 탐구: 물체와 물질 변화 관찰하고 예측하기
"""

# 입력 폼
st.subheader("📝 아이 정보 입력")

with st.form("child_info_form"):
    col1, col2 = st.columns(2)

    with col1:
        child_name = st.text_input("아이 이름", placeholder="예: 송기훈")
        child_age = st.selectbox("나이", ["만 3세", "만 4세", "만 5세", "만 6세", "만 7세"])

    with col2:
        semester = st.radio("학기", ["1학기", "2학기"], horizontal=True)

    observations = st.text_area(
        "관찰 내용",
        placeholder=(
            "예: 참신하고, 창의로운 이야기들을 주로함. "
            "경험한 일에 대해 기억하고 교사나 친구들에게 이야기함. "
            "도전하는 음식이 많고 골고루 먹으며 남기지 않고 먹는 모습이 많아짐."
        ),
        height=150
    )

    submitted = st.form_submit_button("✍️ 평가문 생성", use_container_width=True)

# 평가문 생성
if submitted:
    if not child_name or not observations:
        st.error("아이 이름과 관찰 내용을 모두 입력해주세요.")
    else:
        user_message = f"""
아래 정보를 바탕으로 유아발달상황 종합평가를 작성해주세요.

**아이 정보**
- 이름: {child_name}
- 나이: {child_age}
- 학기: {semester}

**관찰 내용**
{observations}

**요청사항**
1. 따뜻하고 상냥한 교사 말투로 작성
2. 신체운동건강, 의사소통, 사회관계, 예술경험, 자연탐구 5개 영역 모두 포함
3. "{child_name}의 2학기 생활도 응원합니다!" 형식으로 긍정적 마무리
4. 가정 연계 지원 방안 제시
5. 자연스러운 문맥 연결
"""

        st.session_state.messages.append({
            "role": "user",
            "content": user_message,
            "timestamp": datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        })

        with st.spinner("🔍 관련 자료를 검색하고 있습니다..."):
            relevant_docs = []
            rag_context = DEFAULT_KNOWLEDGE

            if mode == "고급 (문서 업로드 포함)" and st.session_state.rag_documents:
                query = f"{child_age} {observations}"
                relevant_docs = search_similar_documents(query, st.session_state.rag_documents, client, top_k=3)

                if relevant_docs:
                    rag_context += "\n\n# 업로드된 참고 자료\n"
                    for i, doc in enumerate(relevant_docs):
                        rag_context += f"\n## 참고 문서 {i + 1}: {doc['filename']} (조각 {doc['chunk_id']})\n"
                        rag_context += doc["content"][:1000] + "\n"

        st.session_state.last_rag_context = rag_context

        with st.spinner("🤖 AI가 평가문을 작성하고 있습니다..."):
            try:
                response = client.chat.completions.create(
                    model="gpt-4o",
                    messages=[
                        {"role": "system", "content": SYSTEM_PROMPT},
                        {"role": "system", "content": f"참고 자료:\n{rag_context}"},
                        {"role": "user", "content": user_message}
                    ],
                    temperature=0.7,
                    max_tokens=2500
                )

                assistant_message = response.choices[0].message.content

                st.session_state.messages.append({
                    "role": "assistant",
                    "content": assistant_message,
                    "timestamp": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                    "relevant_docs": relevant_docs
                })

            except Exception as exc:
                st.error(f"오류가 발생했습니다: {str(exc)}")

# 대화 내역 표시
if st.session_state.messages:
    st.markdown("---")
    st.subheader("💬 평가문 및 대화 내역")

    for i, message in enumerate(st.session_state.messages):
        if message["role"] == "user":
            with st.chat_message("user", avatar="👤"):
                st.markdown(message["content"])
                st.caption(f"📅 {message['timestamp']}")
        else:
            with st.chat_message("assistant", avatar="🤖"):
                content = message["content"]

                if "---" in content and ("**📚 참고한" in content or "**🔍 발달적" in content):
                    parts = content.split("---")
                    evaluation = parts[0].strip()
                    reference = "---".join(parts[1:]).strip() if len(parts) > 1 else ""

                    st.markdown("### 📄 작성된 평가문")
                    st.markdown(evaluation)

                    if reference:
                        with st.expander("🔍 참고 자료 및 발달적 근거 보기"):
                            st.markdown(reference)

                            if "relevant_docs" in message and message["relevant_docs"]:
                                st.markdown("---")
                                st.markdown("### 📚 검색된 유사 문서")
                                for j, doc in enumerate(message["relevant_docs"]):
                                    with st.expander(f"문서 {j + 1}: {doc['filename']}"):
                                        st.text(doc["content"][:500] + "...")
                else:
                    st.markdown(content)

                st.caption(f"📅 {message['timestamp']}")

                col1, col2, _ = st.columns([1, 1, 4])
                with col1:
                    if st.button("📋 복사", key=f"copy_{i}"):
                        st.code(content, language=None)
                with col2:
                    if st.button("💾 저장", key=f"save_{i}"):
                        SAVE_DIR.mkdir(exist_ok=True)
                        filename = f"{datetime.now().strftime('%Y%m%d_%H%M%S')}.txt"
                        with open(SAVE_DIR / filename, "w", encoding="utf-8") as f:
                            f.write(content)
                        st.success(f"✅ {filename}으로 저장되었습니다!")

    st.markdown("---")
    follow_up = st.text_input(
        "💭 수정 요청이나 추가 질문이 있으신가요?",
        placeholder="예: 사회관계 부분을 더 구체적으로 써주세요"
    )

    if st.button("📤 전송") and follow_up:
        st.session_state.messages.append({
            "role": "user",
            "content": follow_up,
            "timestamp": datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        })

        with st.spinner("🤖 답변을 생성하고 있습니다..."):
            try:
                rag_context = st.session_state.last_rag_context or DEFAULT_KNOWLEDGE
                conversation_history = [
                    {"role": "system", "content": SYSTEM_PROMPT},
                    {"role": "system", "content": f"참고 자료:\n{rag_context}"}
                ]

                for msg in st.session_state.messages:
                    conversation_history.append({
                        "role": msg["role"],
                        "content": msg["content"]
                    })

                response = client.chat.completions.create(
                    model="gpt-4o",
                    messages=conversation_history,
                    temperature=0.7,
                    max_tokens=2500
                )

                assistant_message = response.choices[0].message.content

                st.session_state.messages.append({
                    "role": "assistant",
                    "content": assistant_message,
                    "timestamp": datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                })

                st.rerun()

            except Exception as exc:
                st.error(f"오류가 발생했습니다: {str(exc)}")

# 푸터
st.markdown("---")
st.markdown(
    """
<div style='text-align: center; color: #666; padding: 20px;'>
    <p>🎓 유아발달상황 종합평가 작성 도우미 (기본/고급 통합)</p>
    <p style='font-size: 0.9em;'>누리과정 5개 영역(신체운동건강, 의사소통, 사회관계, 예술경험, 자연탐구) 기반</p>
    <p style='font-size: 0.8em; color: #999;'>OpenAI GPT-4 & Embeddings API 사용</p>
</div>
""",
    unsafe_allow_html=True
)
